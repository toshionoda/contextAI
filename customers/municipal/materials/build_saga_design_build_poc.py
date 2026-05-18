"""
佐賀市 設計施工POC実証事業（ご提案）PowerPoint生成スクリプト

既存の『佐賀市DX導入提案r2.pptx』のデザインを踏襲:
- サイズ: 10.00 x 5.62 in
- フォント: Meiryo
- メイン色: #1B3A5C / アクセント: #2E7D32
- カード背景: #F5F7FA, #EFF5EF / 本文: #333333
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ===== 色 =====
NAVY = RGBColor(0x1B, 0x3A, 0x5C)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
LIGHT_BLUE = RGBColor(0xB0, 0xC4, 0xDE)
CARD_GRAY = RGBColor(0xF5, 0xF7, 0xFA)
CARD_GREEN = RGBColor(0xEF, 0xF5, 0xEF)
BODY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DIVIDER = RGBColor(0xCC, 0xCC, 0xCC)
RED = RGBColor(0xC6, 0x28, 0x28)
FONT = 'Meiryo'

SLIDE_W = 10.0
SLIDE_H = 5.62

prs = Presentation()
prs.slide_width = Inches(SLIDE_W)
prs.slide_height = Inches(SLIDE_H)

blank = prs.slide_layouts[6]  # blank


def set_text(frame, text, size=10, bold=False, color=BODY, font=FONT, align=None, anchor=None):
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.text = ''
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if align is not None:
        p.alignment = align
    if anchor is not None:
        frame.vertical_anchor = anchor
    # Clear default margins
    frame.margin_left = Emu(36000)
    frame.margin_right = Emu(36000)
    frame.margin_top = Emu(18000)
    frame.margin_bottom = Emu(18000)


def add_text(slide, x, y, w, h, text, size=10, bold=False, color=BODY, align=None, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(tb.text_frame, text, size=size, bold=bold, color=color, align=align, anchor=anchor)
    return tb


def add_multi_text(slide, x, y, w, h, runs, anchor=MSO_ANCHOR.TOP, align=None, line_spacing=None):
    """runs: list of dicts {text, size, bold, color, align}"""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(36000)
    tf.margin_right = Emu(36000)
    tf.margin_top = Emu(18000)
    tf.margin_bottom = Emu(18000)
    tf.vertical_anchor = anchor
    first = True
    for r in runs:
        if r.get('newline'):
            p = tf.add_paragraph()
            first = False
            continue
        if r.get('newpara'):
            p = tf.add_paragraph()
            if align is not None:
                p.alignment = align
            if r.get('align') is not None:
                p.alignment = r['align']
            if line_spacing:
                p.line_spacing = line_spacing
            first = False
            continue
        if first:
            p = tf.paragraphs[0]
            p.text = ''
            if align is not None:
                p.alignment = align
            first = False
        run = p.add_run()
        run.text = r['text']
        run.font.name = r.get('font', FONT)
        run.font.size = Pt(r.get('size', 10))
        run.font.bold = r.get('bold', False)
        run.font.color.rgb = r.get('color', BODY)
    return tb


def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.shadow.inherit = False
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        if line_w is not None:
            shp.line.width = Pt(line_w)
    shp.text_frame.margin_left = Emu(36000)
    shp.text_frame.margin_right = Emu(36000)
    shp.text_frame.margin_top = Emu(18000)
    shp.text_frame.margin_bottom = Emu(18000)
    return shp


def add_line(slide, x1, y1, x2, y2, color=DIVIDER, width=0.75):
    ln = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(width)
    return ln


def add_header(slide, title, sub_right=None):
    """各コンテンツスライド共通のヘッダ"""
    # タイトル
    add_text(slide, 0.42, 0.17, 9.0, 0.47, title, size=22, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    # 水平線
    add_line(slide, 0.0, 0.72, SLIDE_W, 0.72, color=DIVIDER, width=0.75)
    if sub_right:
        add_text(slide, 7.5, 0.25, 2.3, 0.35, sub_right, size=9, color=RGBColor(0x80, 0x80, 0x80), align=PP_ALIGN.RIGHT)


def set_slide_bg(slide, color):
    """スライド背景を単色で塗る"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# ==================================================================
# Slide 1: 表紙
# ==================================================================
s1 = prs.slides.add_slide(blank)
set_slide_bg(s1, NAVY)
# Main title
add_text(s1, 0.77, 1.60, 8.46, 0.80,
         '佐賀市 設計施工POC実証事業（ご提案）',
         size=30, bold=True, color=WHITE)
# Subtitle
add_text(s1, 0.79, 2.55, 8.42, 0.40,
         '〜 落札企業の設計施工業務を起点に CiviLink で業務効率化を実証する 〜',
         size=15, color=LIGHT_BLUE)
# Divider
add_line(s1, 3.5, 3.20, 6.5, 3.20, color=LIGHT_BLUE, width=0.75)
# Date
add_text(s1, 3.50, 3.35, 3.0, 0.28, '令和8年4月', size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)
add_text(s1, 3.50, 3.70, 3.0, 0.28, '株式会社マルメ', size=13, color=LIGHT_BLUE, align=PP_ALIGN.CENTER)


# ==================================================================
# Slide 2: 本POCの位置づけ
# ==================================================================
s2 = prs.slides.add_slide(blank)
add_header(s2, '本POCの位置づけ')

# リード文
add_multi_text(s2, 0.42, 0.82, 9.17, 0.60, [
    {'text': '佐賀市が公募予定の投資提案事業を起点に、', 'size': 11, 'color': BODY},
    {'text': '落札企業（地場建設会社）の設計施工業務', 'size': 11, 'bold': True, 'color': NAVY},
    {'text': 'に入る段階でCiviLinkを導入。', 'size': 11, 'color': BODY},
    {'newpara': True},
    {'text': '実工事の中で「どれだけ業務が効率化されるか」を定量的に実証し、包括委託への足場とする。', 'size': 11, 'color': BODY},
])

# タイムライン矢印
# 4ステップ：投資提案→落札→設計施工（★POC）→包括委託
steps = [
    ('STEP 1', '2026年5月〜', '投資提案の公募', '佐賀市が地域建設業DX\n普及促進事業を公募', CARD_GRAY, NAVY),
    ('STEP 2', '〜2026年8月', '3〜4ヶ月で落札', '地場建設会社\n（Cランク相当）が受託', CARD_GRAY, NAVY),
    ('STEP 3', '2026年9月〜\n2027年3月', '★ 本POC実施', '設計施工業務にCiviLinkを\n導入し効率化を実証', CARD_GREEN, GREEN),
    ('STEP 4', '2027年4月〜', '包括委託へ接続', '実証成果を本委託の\nDX標準仕様に反映', CARD_GRAY, NAVY),
]

card_w = 2.18
card_h = 2.80
gap = 0.20
total_w = card_w * 4 + gap * 3
start_x = (SLIDE_W - total_w) / 2
card_y = 1.72

for i, (step, period, title_txt, desc, bg, accent) in enumerate(steps):
    x = start_x + i * (card_w + gap)
    add_rect(s2, x, card_y, card_w, card_h, fill=bg)
    # 上部 色バンド
    add_rect(s2, x, card_y, card_w, 0.28, fill=accent)
    add_text(s2, x, card_y, card_w, 0.28, step, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 期間
    add_text(s2, x + 0.05, card_y + 0.35, card_w - 0.10, 0.50, period, size=9, bold=True, color=accent, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    # タイトル
    add_text(s2, x + 0.05, card_y + 0.90, card_w - 0.10, 0.55, title_txt, size=12, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)
    # 説明
    add_text(s2, x + 0.10, card_y + 1.55, card_w - 0.20, 1.15, desc, size=9, color=BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# 矢印記号
for i in range(3):
    ax = start_x + card_w + (card_w + gap) * i - 0.10
    add_text(s2, ax, card_y + 1.20, 0.20, 0.40, '▶', size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# メッセージ
add_rect(s2, 0.42, 4.75, 9.17, 0.65, fill=CARD_GRAY)
add_multi_text(s2, 0.55, 4.80, 8.90, 0.55, [
    {'text': 'POC実施のねらい： ', 'size': 11, 'bold': True, 'color': NAVY},
    {'text': '「実際に発注する工事」で効果を測定し、包括委託のDX標準を佐賀市発で確立する。', 'size': 11, 'color': BODY},
], anchor=MSO_ANCHOR.MIDDLE)


# ==================================================================
# Slide 3: POCの目的
# ==================================================================
s3 = prs.slides.add_slide(blank)
add_header(s3, 'POCの目的')

add_text(s3, 0.42, 0.82, 9.17, 0.30,
         '設計施工フェーズに CiviLink を導入することで、3つの価値を同時に実証する。',
         size=11, color=BODY)

# 3カード
targets = [
    ('① 落札企業の働き方改革', '地場建設会社（Cランク）の\n設計施工業務を効率化', [
        '紙・FAX・押印待ちの解消',
        '書類作成工数の削減',
        '少人数でも工事を回せる体制の構築',
    ]),
    ('② 市側監督業務の効率化', '佐賀市の発注者監督業務の\n負担を軽減', [
        '照査・承認のリードタイム短縮',
        '指摘履歴・判断ログのデジタル蓄積',
        '電子納品の形式エラー削減',
    ]),
    ('③ 3者連携モデルの実地検証', '自治体×補修設計×建設の\n一気通貫運用を検証', [
        '同一図面を軸にしたデータ連携',
        '包括委託の標準仕様化に向けた知見',
        '他地場企業・他工事への横展開',
    ]),
]
cw = 3.00
ch = 3.50
cgap = 0.13
tx = (SLIDE_W - (cw * 3 + cgap * 2)) / 2
cy = 1.30
for i, (ttl, sub, bullets) in enumerate(targets):
    x = tx + i * (cw + cgap)
    add_rect(s3, x, cy, cw, ch, fill=CARD_GREEN)
    # 色バンド
    add_rect(s3, x, cy, cw, 0.05, fill=GREEN)
    add_text(s3, x + 0.12, cy + 0.15, cw - 0.24, 0.30, ttl, size=12, bold=True, color=GREEN, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s3, x + 0.12, cy + 0.50, cw - 0.24, 0.70, sub, size=11, bold=True, color=NAVY)
    by = cy + 1.35
    for b in bullets:
        add_multi_text(s3, x + 0.12, by, cw - 0.24, 0.55, [
            {'text': '● ', 'size': 9, 'color': GREEN, 'bold': True},
            {'text': b, 'size': 9.5, 'color': BODY},
        ])
        by += 0.62

# フッターメッセージ
add_rect(s3, 0.42, 5.00, 9.17, 0.45, fill=NAVY)
add_text(s3, 0.55, 5.00, 8.90, 0.45,
         '得られた効果・知見は、佐賀市×伊藤忠の包括委託（2027年度〜）のDX標準仕様に直結する。',
         size=10.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ==================================================================
# Slide 4: 実証対象・前提
# ==================================================================
s4 = prs.slides.add_slide(blank)
add_header(s4, '実証対象・前提')

# 左: 対象サマリテーブル
add_text(s4, 0.42, 0.82, 4.58, 0.30, '実証対象サマリ', size=12, bold=True, color=NAVY)

rows = [
    ('発注者', '佐賀市'),
    ('落札企業', 'Cランク相当 地場建設会社（1〜2社）'),
    ('補修設計会社', '既存の契約設計会社（市側で選定）'),
    ('対象工事', '佐賀市発注の補修工事（橋梁・道路等）'),
    ('実証期間', '2026年9月〜2027年3月（約7ヶ月）'),
    ('活用ツール', 'CiviLink + DXパッケージ（写真AI・AI点検）'),
    ('実証費用', '落札工事費に加算せず、マルメ＋市で別途検討'),
]
tbl_y = 1.20
row_h = 0.40
for i, (k, v) in enumerate(rows):
    y = tbl_y + i * row_h
    bg = CARD_GRAY if i % 2 == 0 else WHITE
    add_rect(s4, 0.42, y, 4.49, row_h, fill=bg)
    add_text(s4, 0.52, y, 1.30, row_h, k, size=10, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s4, 1.85, y, 3.00, row_h, v, size=10, color=BODY, anchor=MSO_ANCHOR.MIDDLE)

# 右: Cランク企業の想定プロフィール
add_text(s4, 5.10, 0.82, 4.58, 0.30, '対象企業（Cランク相当）の想定プロフィール', size=12, bold=True, color=NAVY)
add_rect(s4, 5.10, 1.20, 4.49, 3.80, fill=CARD_GRAY)

profile_items = [
    ('従業員規模', '技術者5〜20名程度の地場建設会社'),
    ('受注工事', '佐賀市発注の補修・維持工事が中心'),
    ('DX状況', 'CAD・Excelが中心。クラウド運用は限定的'),
    ('課題', '書類作成の属人化、若手確保、2024年問題'),
    ('DX適性', '補助金・NETIS加点に関心あり。伴走支援を求める'),
]
py = 1.35
for k, v in profile_items:
    add_multi_text(s4, 5.25, py, 4.30, 0.70, [
        {'text': k, 'size': 10, 'bold': True, 'color': GREEN},
        {'newpara': True},
        {'text': v, 'size': 10, 'color': BODY},
    ])
    py += 0.72

# 下部補足
add_rect(s4, 0.42, 5.10, 9.17, 0.35, fill=CARD_GREEN)
add_text(s4, 0.55, 5.10, 8.90, 0.35,
         '※ 対象工事・対象企業は佐賀市と協議の上で選定。補修設計会社の参画も含めた3者連携で実施する。',
         size=9.5, color=BODY, anchor=MSO_ANCHOR.MIDDLE)


# ==================================================================
# Slide 5: 現状の業務課題
# ==================================================================
s5 = prs.slides.add_slide(blank)
add_header(s5, '現状の設計施工業務の課題')

add_text(s5, 0.42, 0.82, 9.17, 0.30,
         '地場建設会社の設計施工業務では、各フェーズで紙・メール・対面中心の運用が残り、手戻りと属人化を招いている。',
         size=10.5, color=BODY)

# 3フェーズのカード
phases = [
    ('① 設計・照査フェーズ', [
        ('紙図面・PDFメール送付', '版管理が追えず最新版が分からない'),
        ('数量拾いの手計算', '図面⇔数量計算書の不整合が発生'),
        ('照査指摘が紙・メール', '指摘履歴が残らず同じ手戻りを繰り返す'),
    ]),
    ('② 施工フェーズ', [
        ('打合簿が紙・押印運用', '監督員との承認に数日〜1週間かかる'),
        ('写真整理・台帳作成が手作業', '現場1工事で数日の事務工数が発生'),
        ('段階確認の連絡が電話・FAX', '記録が残らず監督履歴が不透明'),
    ]),
    ('③ 検査・電子納品フェーズ', [
        ('納品フォルダを手作業で構成', '形式エラーで差戻しが多発'),
        ('成果物が各社バラバラの構成', '市側のチェック負担が大きい'),
        ('過去成果物の横断参照が困難', '類似工事の知見が再利用されない'),
    ]),
]
cw = 3.03
ch = 3.80
cgap = 0.11
tx = (SLIDE_W - (cw * 3 + cgap * 2)) / 2
cy = 1.20
for i, (ttl, items) in enumerate(phases):
    x = tx + i * (cw + cgap)
    add_rect(s5, x, cy, cw, ch, fill=CARD_GRAY)
    add_rect(s5, x, cy, cw, 0.40, fill=NAVY)
    add_text(s5, x, cy, cw, 0.40, ttl, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    by = cy + 0.55
    for fact, impact in items:
        add_multi_text(s5, x + 0.12, by, cw - 0.24, 1.05, [
            {'text': '● ' + fact, 'size': 10, 'bold': True, 'color': NAVY},
            {'newpara': True},
            {'text': '    → ' + impact, 'size': 9, 'color': BODY},
        ])
        by += 1.08

# 下部メッセージ
add_rect(s5, 0.42, 5.10, 9.17, 0.35, fill=NAVY)
add_text(s5, 0.55, 5.10, 8.90, 0.35,
         '各フェーズの「紙・対面・属人化」をCiviLinkで「クラウド・ワークフロー・標準化」に置き換える。',
         size=10.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ==================================================================
# Slide 6: Before/After 業務フロー
# ==================================================================
s6 = prs.slides.add_slide(blank)
add_header(s6, 'CiviLink導入による業務フローの変化')

# フロー5ステップ：受領 → 設計・照査 → 施工準備 → 施工・検査 → 電子納品
flow = ['図面・設計図書の受領', '設計・照査', '施工準備・着工', '施工・段階確認', '検査・電子納品']

# 共通レイアウト
step_w = 1.76
step_gap = 0.12
total_w = step_w * 5 + step_gap * 4
flow_x = (SLIDE_W - total_w) / 2

# Before 行
add_text(s6, 0.42, 0.88, 1.40, 0.30, 'Before', size=12, bold=True, color=RED)
add_text(s6, 1.70, 0.88, 6.0, 0.30, '紙・メール・押印・手作業', size=10, color=BODY)

before_y = 1.22
before_h = 1.40
for i, label in enumerate(flow):
    x = flow_x + i * (step_w + step_gap)
    add_rect(s6, x, before_y, step_w, 0.35, fill=RED)
    add_text(s6, x, before_y, step_w, 0.35, label, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s6, x, before_y + 0.35, step_w, before_h - 0.35, fill=CARD_GRAY)

before_contents = [
    '紙図面・PDFメール\n版管理なし',
    '個人PCで照査\n指摘は紙・口頭',
    '施工計画書\nゼロから作成',
    '打合簿・写真\n紙・手作業整理',
    '納品フォルダ\n手作業構成',
]
for i, txt in enumerate(before_contents):
    x = flow_x + i * (step_w + step_gap)
    add_text(s6, x + 0.08, before_y + 0.40, step_w - 0.16, before_h - 0.42, txt, size=9, color=BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 矢印 ▼
add_text(s6, 0.42, 2.70, 9.17, 0.35, '▼  CiviLink 導入  ▼', size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# After 行
add_text(s6, 0.42, 3.15, 1.40, 0.30, 'After', size=12, bold=True, color=GREEN)
add_text(s6, 1.70, 3.15, 6.0, 0.30, 'クラウド・ワークフロー・AI支援', size=10, color=BODY)

after_y = 3.48
after_h = 1.55
for i, label in enumerate(flow):
    x = flow_x + i * (step_w + step_gap)
    add_rect(s6, x, after_y, step_w, 0.35, fill=GREEN)
    add_text(s6, x, after_y, step_w, 0.35, label, size=9.5, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s6, x, after_y + 0.35, step_w, after_h - 0.35, fill=CARD_GREEN)

after_contents = [
    'クラウド図面\n版管理・共有',
    'デジタル照査\n指摘管理＋AI',
    '施工計画AIで\n初稿30分生成',
    '打合簿クラウド\n承認WF/写真AI',
    '電子納品を\n自動構成',
]
for i, txt in enumerate(after_contents):
    x = flow_x + i * (step_w + step_gap)
    add_text(s6, x + 0.08, after_y + 0.40, step_w - 0.16, after_h - 0.42, txt, size=9, color=BODY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# 下部メッセージ
add_text(s6, 0.42, 5.15, 9.17, 0.30,
         '各フェーズの情報が同一のCiviLink上に蓄積され、受発注者間の「探す・待つ・作り直す」がなくなる。',
         size=10.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ==================================================================
# Slide 7: 設計フェーズの実証内容
# ==================================================================
s7 = prs.slides.add_slide(blank)
add_header(s7, '設計フェーズの実証内容')

add_text(s7, 0.42, 0.82, 9.17, 0.30,
         '補修設計成果物の受領〜照査〜施工への受渡しまでを、クラウド上で完結させる。',
         size=10.5, color=BODY)

design_items = [
    ('① 図面・設計図書のクラウド共有', 'CiviLink図面管理', [
        '補修設計会社 → 建設会社 → 市への三者共有',
        '常に最新版が全員に可視化（版管理）',
        'スマホ・タブレットから現場で閲覧',
    ], '図面受渡しメール往復の削減'),
    ('② デジタル照査・指摘管理', 'CiviLink指摘管理WF', [
        '図面上に直接コメント・指摘を記録',
        '指摘→回答→確認の履歴が残る',
        '過去工事の指摘パターンを横断検索',
    ], '照査時間短縮／品質ばらつき解消'),
    ('③ 図面⇔数量の整合チェック', 'CiviLink AI照査', [
        '図面と数量計算書の不整合をAIが検出',
        '桁違い・単位ミスを自動ハイライト',
        '補修設計会社の自己照査にも活用',
    ], '積算ミス／手戻りの削減'),
]
cw = 3.03
ch = 3.90
cgap = 0.11
tx = (SLIDE_W - (cw * 3 + cgap * 2)) / 2
cy = 1.20
for i, (ttl, tool, bullets, effect) in enumerate(design_items):
    x = tx + i * (cw + cgap)
    add_rect(s7, x, cy, cw, ch, fill=CARD_GREEN)
    add_rect(s7, x, cy, cw, 0.42, fill=GREEN)
    add_text(s7, x, cy, cw, 0.42, ttl, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # 活用機能
    add_rect(s7, x + 0.12, cy + 0.55, cw - 0.24, 0.32, fill=WHITE)
    add_text(s7, x + 0.12, cy + 0.55, cw - 0.24, 0.32, '活用機能：' + tool, size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # bullets
    by = cy + 1.00
    for b in bullets:
        add_multi_text(s7, x + 0.15, by, cw - 0.30, 0.65, [
            {'text': '● ', 'size': 9, 'color': GREEN, 'bold': True},
            {'text': b, 'size': 9.5, 'color': BODY},
        ])
        by += 0.68
    # effect
    add_rect(s7, x + 0.12, cy + ch - 0.70, cw - 0.24, 0.58, fill=NAVY)
    add_multi_text(s7, x + 0.18, cy + ch - 0.70, cw - 0.36, 0.58, [
        {'text': '検証する効果', 'size': 8.5, 'bold': True, 'color': LIGHT_BLUE},
        {'newpara': True},
        {'text': effect, 'size': 9.5, 'bold': True, 'color': WHITE},
    ], anchor=MSO_ANCHOR.MIDDLE)

add_text(s7, 0.42, 5.20, 9.17, 0.30,
         '※ 検証では工数・リードタイム・指摘件数を導入前後で比較し、定量効果を測定する。',
         size=9.5, color=BODY, anchor=MSO_ANCHOR.MIDDLE)


# ==================================================================
# Slide 8: 施工フェーズの実証内容
# ==================================================================
s8 = prs.slides.add_slide(blank)
add_header(s8, '施工フェーズの実証内容')

add_text(s8, 0.42, 0.82, 9.17, 0.30,
         '着工〜段階確認〜検査〜電子納品までを、打合簿・写真・承認のクラウド化で効率化する。',
         size=10.5, color=BODY)

const_items = [
    ('① 打合簿・承認ワークフロー', 'CiviLink承認WF', [
        '施工計画書・変更協議をクラウド提出',
        '監督員の承認・差戻しを即時化',
        '承認履歴・判断ログが自動蓄積',
    ], '承認リードタイム ▲70%目標'),
    ('② 出来形・写真台帳の自動化', 'CiviLink × 写真AI連携', [
        '現場写真を黒板AI-OCRで自動分類',
        '出来形測定値を現場からクラウド入力',
        '写真台帳が自動生成（手作業排除）',
    ], '台帳作成工数 ▲90%目標'),
    ('③ 電子納品の自動構成', 'CiviLink電子納品', [
        '日々の書類・写真から納品フォルダを自動構成',
        '国交省基準への適合性を自動チェック',
        '市側の差戻し発生ゼロを目標',
    ], '納品差戻し件数 ゼロ目標'),
]
cw = 3.03
ch = 3.90
cgap = 0.11
tx = (SLIDE_W - (cw * 3 + cgap * 2)) / 2
cy = 1.20
for i, (ttl, tool, bullets, effect) in enumerate(const_items):
    x = tx + i * (cw + cgap)
    add_rect(s8, x, cy, cw, ch, fill=CARD_GREEN)
    add_rect(s8, x, cy, cw, 0.42, fill=GREEN)
    add_text(s8, x, cy, cw, 0.42, ttl, size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_rect(s8, x + 0.12, cy + 0.55, cw - 0.24, 0.32, fill=WHITE)
    add_text(s8, x + 0.12, cy + 0.55, cw - 0.24, 0.32, '活用機能：' + tool, size=9.5, bold=True, color=NAVY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    by = cy + 1.00
    for b in bullets:
        add_multi_text(s8, x + 0.15, by, cw - 0.30, 0.65, [
            {'text': '● ', 'size': 9, 'color': GREEN, 'bold': True},
            {'text': b, 'size': 9.5, 'color': BODY},
        ])
        by += 0.68
    add_rect(s8, x + 0.12, cy + ch - 0.70, cw - 0.24, 0.58, fill=NAVY)
    add_multi_text(s8, x + 0.18, cy + ch - 0.70, cw - 0.36, 0.58, [
        {'text': '検証する効果', 'size': 8.5, 'bold': True, 'color': LIGHT_BLUE},
        {'newpara': True},
        {'text': effect, 'size': 9.5, 'bold': True, 'color': WHITE},
    ], anchor=MSO_ANCHOR.MIDDLE)

add_text(s8, 0.42, 5.20, 9.17, 0.30,
         '※ 各KPIは実工事の業務データから測定。測定結果は市・建設会社・マルメで共有する。',
         size=9.5, color=BODY, anchor=MSO_ANCHOR.MIDDLE)


# ==================================================================
# Slide 9: 期待効果（KPI）
# ==================================================================
s9 = prs.slides.add_slide(blank)
add_header(s9, '期待効果（KPI）')

add_text(s9, 0.42, 0.82, 9.17, 0.30,
         '設計・施工の各業務で、定量的な効果目標を設定して検証する。',
         size=10.5, color=BODY)

# テーブルヘッダ
headers = ['フェーズ', '業務領域', '現状', '目標値', '測定方法']
widths = [0.95, 2.05, 2.20, 2.25, 1.72]
ty = 1.20
tx = 0.42
header_h = 0.42

hx = tx
for i, h in enumerate(headers):
    add_rect(s9, hx, ty, widths[i], header_h, fill=NAVY)
    add_text(s9, hx, ty, widths[i], header_h, h, size=10, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    hx += widths[i]

kpi_rows = [
    ('設計', '図面⇔数量整合チェック', '手作業・属人対応', 'AI自動検出で照査工数▲50%', '照査時間の工数記録'),
    ('設計', '照査指摘のやり取り', '紙・メール往復', '平均回答時間 ▲60%', '指摘ログのタイムスタンプ'),
    ('施工', '打合簿・承認', '押印・対面で3〜7日', '承認リードタイム ▲70%', '承認履歴の日数集計'),
    ('施工', '写真台帳・出来形', '1工事あたり数日', '台帳作成工数 ▲90%', '作業工数の前後比較'),
    ('施工', '段階確認の記録', '電話・FAX中心', 'デジタル化率 100%', '記録媒体の内訳'),
    ('納品', '電子納品', '形式エラー差戻し多発', '差戻し件数 ゼロ', '市側検査結果'),
]
row_h = 0.38
for i, row in enumerate(kpi_rows):
    y = ty + header_h + i * row_h
    hx = tx
    bg = CARD_GRAY if i % 2 == 0 else WHITE
    for j, val in enumerate(row):
        add_rect(s9, hx, y, widths[j], row_h, fill=bg)
        if j == 0:
            add_text(s9, hx, y, widths[j], row_h, val, size=9.5, bold=True, color=GREEN, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        elif j == 3:
            add_text(s9, hx + 0.08, y, widths[j] - 0.16, row_h, val, size=9.5, bold=True, color=NAVY, anchor=MSO_ANCHOR.MIDDLE)
        else:
            add_text(s9, hx + 0.08, y, widths[j] - 0.16, row_h, val, size=9, color=BODY, anchor=MSO_ANCHOR.MIDDLE)
        hx += widths[j]

# 下部メッセージ
add_rect(s9, 0.42, 5.00, 9.17, 0.45, fill=CARD_GREEN)
add_multi_text(s9, 0.55, 5.00, 8.90, 0.45, [
    {'text': 'POC成果の活用： ', 'size': 10.5, 'bold': True, 'color': GREEN},
    {'text': '実測値をもとに、包括委託でのDX標準仕様・費用効果指標を策定する。', 'size': 10.5, 'color': BODY},
], anchor=MSO_ANCHOR.MIDDLE)


# ==================================================================
# Slide 10: 実施体制・スケジュール
# ==================================================================
s10 = prs.slides.add_slide(blank)
add_header(s10, '実施体制・スケジュール')

# 左: 役割分担
add_text(s10, 0.42, 0.82, 4.58, 0.30, '役割分担', size=12, bold=True, color=NAVY)

roles = [
    ('佐賀市', '対象工事選定／監督員参加／成果評価', NAVY),
    ('落札企業\n（Cランク地場建設）', 'CiviLink利用／業務データ提供／\n効果測定への協力', GREEN),
    ('補修設計会社', '設計成果物の提供／3者連携運用への参画', NAVY),
    ('マルメ', 'CiviLink提供／導入支援／効果測定／\n伴走コンサルティング', GREEN),
]
ry = 1.20
for name, desc, acc in roles:
    add_rect(s10, 0.42, ry, 4.49, 0.78, fill=CARD_GRAY)
    add_rect(s10, 0.42, ry, 0.10, 0.78, fill=acc)
    add_text(s10, 0.62, ry, 1.60, 0.78, name, size=10.5, bold=True, color=acc, anchor=MSO_ANCHOR.MIDDLE)
    add_text(s10, 2.25, ry, 2.55, 0.78, desc, size=9.5, color=BODY, anchor=MSO_ANCHOR.MIDDLE)
    ry += 0.86

# 右: スケジュール
add_text(s10, 5.10, 0.82, 4.58, 0.30, '実施スケジュール', size=12, bold=True, color=NAVY)

schedule = [
    ('2026年9〜10月', '準備フェーズ', '対象工事・対象企業の選定／CiviLink環境構築／キックオフ'),
    ('2026年10〜12月', '設計フェーズ実証', '図面共有・デジタル照査・数量整合チェックの運用'),
    ('2027年1〜3月', '施工フェーズ実証', '打合簿・承認WF・写真台帳自動化・電子納品の運用'),
    ('2027年3月', '効果測定・報告', 'KPI集計／成果報告／包括委託への反映事項整理'),
]
sy = 1.20
for period, phase, desc in schedule:
    add_rect(s10, 5.10, sy, 4.49, 0.90, fill=CARD_GRAY)
    add_rect(s10, 5.10, sy, 0.10, 0.90, fill=GREEN)
    add_text(s10, 5.28, sy + 0.05, 4.20, 0.22, period, size=9.5, bold=True, color=GREEN)
    add_text(s10, 5.28, sy + 0.26, 4.20, 0.28, phase, size=11, bold=True, color=NAVY)
    add_text(s10, 5.28, sy + 0.56, 4.20, 0.32, desc, size=9, color=BODY)
    sy += 0.97

# 下部メッセージ
add_rect(s10, 0.42, 5.10, 9.17, 0.35, fill=NAVY)
add_text(s10, 0.55, 5.10, 8.90, 0.35,
         'まずは対象工事の選定と体制構築から着手。マルメが伴走し、落札企業の負担を最小化する。',
         size=10.5, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)


# ==================================================================
# 保存
# ==================================================================
OUT = 'customers/municipal/materials/saga_design_build_poc_proposal.pptx'
prs.save(OUT)
print(f'Saved: {OUT}')
print(f'Slides: {len(prs.slides)}')
